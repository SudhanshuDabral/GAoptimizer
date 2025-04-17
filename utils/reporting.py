import os
import logging
from logging.handlers import RotatingFileHandler
from io import BytesIO
import time
import plotly.io as pio
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
import pandas as pd
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Logger setup
log_dir = 'logs'
if not os.path.exists(log_dir):
    os.makedirs(log_dir)

log_file = os.path.join(log_dir, 'app_logs.log')

if not logging.getLogger().handlers:
    logging.basicConfig(level=logging.INFO,
                        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
                        handlers=[
                            RotatingFileHandler(log_file, maxBytes=10485760, backupCount=5),
                        ])

logger = logging.getLogger(__name__)

def log_message(level, message, exc_info=False):
    logger.log(level, f"[Reporting] {message}", exc_info=exc_info)

# Configure Kaleido settings
pio.kaleido.scope.default_width = 1600
pio.kaleido.scope.default_height = 900
pio.kaleido.scope.default_scale = 2

def save_plot_as_image(fig, plot_name, user_id):
    """
    Convert a plotly figure to a high-resolution image and save it to a file.
    
    :param fig: plotly Figure object
    :param plot_name: str, name of the plot
    :param user_id: str, user identifier
    :return: str, path to the saved image
    """
    base_dir = 'reporting'
    user_dir = os.path.join(base_dir, user_id)
    os.makedirs(user_dir, exist_ok=True)
    
    filename = f"{plot_name}_{user_id}.png"
    file_path = os.path.join(user_dir, filename)
    
    # Adjust layout to ensure all labels are visible
    fig.update_layout(
        margin=dict(l=150, r=50, t=50, b=50),
        autosize=False,
        width=1600,
        height=900
    )
    
    # Use the format parameter and specify the engine
    pio.write_image(fig, file_path, format='png', engine="kaleido")
    
    return file_path

def cleanup_images(user_id):
    """
    Delete the temporary image files created for a specific user.
    
    :param user_id: str, user identifier
    """
    base_dir = 'reporting'
    user_dir = os.path.join(base_dir, user_id)
    if os.path.exists(user_dir):
        shutil.rmtree(user_dir)
        log_message(logging.INFO, f"Cleaned up temporary images for user {user_id}")

def generate_pdf_report(custom_equation, baseline_productivity, sensitivity_df, 
                        general_results, attribute_specific_results, user_id):
    overall_start_time = time.time()
    log_message(logging.INFO, "Starting PDF report generation")
    buffer = BytesIO()
    try:
        doc = SimpleDocTemplate(buffer, pagesize=landscape(letter), rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=18)
        
        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(name='Justify', alignment=1))
        
        elements = []
        
        # Title
        elements.append(Paragraph("Sensitivity Analysis Report", styles['Title']))
        elements.append(Spacer(1, 12))
        
        # Custom Equation
        elements.append(Paragraph("Custom Equation:", styles['Heading2']))
        elements.append(Paragraph(custom_equation, styles['Normal']))
        elements.append(Spacer(1, 12))
        
        # Baseline Productivity
        elements.append(Paragraph(f"Baseline Productivity: {baseline_productivity:.4f}", styles['Normal']))
        elements.append(Spacer(1, 12))
        
        # General Sensitivity Analysis Results
        elements.append(Paragraph("General Sensitivity Analysis Results:", styles['Heading2']))
        
        # General Sensitivity Analysis Plots
        plot_info = {
            "Tornado Chart": "The Tornado Chart displays the sensitivity of the model to each input variable. Longer bars indicate higher sensitivity.",
            "Feature Importance Chart": "The Feature Importance Chart shows the relative importance of each feature in the model. Taller bars represent more important features.",
            "Elasticity Analysis": "The Elasticity Analysis chart shows how responsive the model output is to changes in each input variable. Steeper lines indicate higher elasticity."
        }
        
        for title, fig in [("Tornado Chart", general_results['fig_tornado']), 
                           ("Feature Importance Chart", general_results['fig_importance']), 
                           ("Elasticity Analysis", general_results['fig_elasticity'])]:
            plot_start_time = time.time()
            log_message(logging.INFO, f"Starting to add {title} to PDF")
            
            elements.append(Paragraph(title, styles['Heading3']))
            elements.append(Paragraph(plot_info[title], styles['Normal']))
            
            try:
                img_path = save_plot_as_image(fig, title.lower().replace(" ", "_"), user_id)
                img = Image(img_path, width=10*inch, height=5.6*inch)
                elements.append(img)
            except Exception as e:
                log_message(logging.ERROR, f"Error adding {title} to PDF: {str(e)}", exc_info=True)
                elements.append(Paragraph(f"Error including {title}: {str(e)}", styles['Normal']))
            
            elements.append(Spacer(1, 12))
            elements.append(PageBreak())
            
            log_message(logging.INFO, f"Finished adding {title} to PDF in {time.time() - plot_start_time:.2f} seconds")
        
        # Attribute-Specific Sensitivity Results
        elements.append(Paragraph("Attribute-Specific Sensitivity Results:", styles['Heading2']))
        for attr, data in attribute_specific_results.items():
            if attr not in ['influence_df', 'fig_influence']:
                elements.append(Paragraph(f"Sensitivity Analysis for {attr}", styles['Heading3']))
                elements.append(Paragraph(f"Attribute Range: {data['attr_min']:.4f} to {data['attr_max']:.4f}", styles['Normal']))
                elements.append(Paragraph(f"Productivity Range: {data['prod_min']:.4f} to {data['prod_max']:.4f}", styles['Normal']))
                elements.append(Paragraph(f"Productivity Spread: {data['prod_range']:.4f}", styles['Normal']))
                
                try:
                    img_path = save_plot_as_image(data['fig'], f"sensitivity_{attr}", user_id)
                    img = Image(img_path, width=10*inch, height=5.6*inch)
                    elements.append(img)
                except Exception as e:
                    log_message(logging.ERROR, f"Error adding attribute-specific plot for {attr} to PDF: {str(e)}", exc_info=True)
                    elements.append(Paragraph(f"Error including plot for {attr}: {str(e)}", styles['Normal']))
                
                elements.append(Spacer(1, 12))
                elements.append(PageBreak())
        
        # Consolidated Attribute Influence
        if 'fig_influence' in attribute_specific_results:
            elements.append(Paragraph("Consolidated Attribute Influence", styles['Heading2']))
            elements.append(Paragraph("This plot represents the overall influence of each attribute on the model's output. Attributes with larger bars have a more significant impact on the productivity predictions.", styles['Normal']))
            try:
                img_path = save_plot_as_image(attribute_specific_results['fig_influence'], "influence_chart", user_id)
                img = Image(img_path, width=10*inch, height=5.6*inch)
                elements.append(img)
            except Exception as e:
                log_message(logging.ERROR, f"Error adding influence chart to PDF: {str(e)}", exc_info=True)
                elements.append(Paragraph(f"Error including influence chart: {str(e)}", styles['Normal']))
            
            elements.append(Spacer(1, 12))
            
            if 'influence_df' in attribute_specific_results:
                elements.append(Paragraph("Detailed Attribute Influence:", styles['Heading3']))
                influence_df = attribute_specific_results['influence_df']
                influence_data = [influence_df.columns.tolist()] + influence_df.values.tolist()
                table_style = TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica'),
                    ('FONTSIZE', (0, 0), (-1, 0), 8),
                    ('TOPPADDING', (0, 0), (-1, -1), 6),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ])
                influence_table = Table(influence_data)
                influence_table.setStyle(table_style)
                elements.append(influence_table)
        
        # Build the PDF
        log_message(logging.DEBUG, "Building PDF")
        pdf_start_time = time.time()
        doc.build(elements)
        log_message(logging.DEBUG, f"Built PDF in {time.time() - pdf_start_time:.2f} seconds")
        
        # Reset buffer position
        buffer.seek(0)
        
        overall_end_time = time.time()
        log_message(logging.INFO, f"PDF report generated successfully in {overall_end_time - overall_start_time:.2f} seconds")
        return buffer
    except Exception as e:
        log_message(logging.ERROR, f"Error in PDF report generation: {str(e)}", exc_info=True)
        raise
    finally:
        if 'doc' in locals():
            del doc  # Explicitly delete the SimpleDocTemplate object
        cleanup_images(user_id)  # Clean up temporary image files

def generate_ppt_report(custom_equation, baseline_productivity, general_results, attribute_specific_results, user_id):
    overall_start_time = time.time()
    log_message(logging.INFO, "Starting PowerPoint report generation")
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    try:
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        title.text = "Sensitivity Analysis Report"
        
        subtitle_text = (
            "This report presents the results of a sensitivity analysis conducted on a hydraulic fracturing productivity model.\n\n"
            f"Baseline Productivity: {baseline_productivity:.4f}\n"
            "Baseline Productivity represents the model's output when all inputs are at their median values.\n\n"
            "The following slides show:\n"
            "1. Model Equation\n"
            "2. General Sensitivity Analysis Results\n"
            "3. Attribute-Specific Sensitivity Results\n"
            "4. Consolidated Attribute Influence\n\n"
            "These analyses help identify which inputs have the most significant impact on the model's predictions."
        )
        if subtitle:
            subtitle.text = subtitle_text
        else:
            # If there's no subtitle placeholder, add a text box
            left = top = Inches(1)
            width = Inches(14)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = subtitle_text

        # Adjust text size and alignment
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.alignment = PP_ALIGN.LEFT

        # Model Equation slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout with title and content
        title = slide.shapes.title
        title.text = "Model Equation"
        
        # Add space between title and equation
        top = Inches(2)  # Increased from 1 to 2 inches
        left = Inches(0.5)
        width = Inches(15)
        height = Inches(5.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Split the equation into multiple lines
        max_chars_per_line = 100
        equation_lines = []
        current_line = ""
        for word in custom_equation.split():
            if len(current_line) + len(word) + 1 <= max_chars_per_line:
                current_line += " " + word if current_line else word
            else:
                equation_lines.append(current_line)
                current_line = word
        if current_line:
            equation_lines.append(current_line)
        
        # Add each line of the equation to the text frame
        for i, line in enumerate(equation_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)  # Increased font size
            p.alignment = PP_ALIGN.LEFT


        # General Sensitivity Analysis Results
        plot_info = {
            "Tornado Chart": "Sensitivity of the model to each input variable",
            "Feature Importance Chart": "Relative importance of each feature in the model",
            "Elasticity Analysis": "Responsiveness of the model output to changes in each input variable"
        }

        for title, fig in [("Tornado Chart", general_results['fig_tornado']), 
                           ("Feature Importance Chart", general_results['fig_importance']), 
                           ("Elasticity Analysis", general_results['fig_elasticity'])]:
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add a text box for the title
            left = top = width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = title
            
            # Add description
            desc_box = slide.shapes.add_textbox(left, Inches(1.5), width, height)
            desc_tf = desc_box.text_frame
            desc_tf.text = plot_info[title]
            
            img_path = save_plot_as_image(fig, title.lower().replace(" ", "_"), user_id)
            slide.shapes.add_picture(img_path, Inches(1), Inches(2.5), width=Inches(14), height=Inches(5.5))

        # Attribute-Specific Sensitivity Results
        for attr, data in attribute_specific_results.items():
            if attr not in ['influence_df', 'fig_influence']:
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add a text box for the title
                left = top = width = height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = f"Sensitivity Analysis for {attr}"
                
                img_path = save_plot_as_image(data['fig'], f"sensitivity_{attr}", user_id)
                slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(14), height=Inches(6))

        # Consolidated Attribute Influence
        if 'fig_influence' in attribute_specific_results:
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add a text box for the title
            left = top = width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "Consolidated Attribute Influence"
            
            img_path = save_plot_as_image(attribute_specific_results['fig_influence'], "influence_chart", user_id)
            slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(14), height=Inches(6))

        # Save the presentation
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)

        overall_end_time = time.time()
        log_message(logging.INFO, f"PowerPoint report generated successfully in {overall_end_time - overall_start_time:.2f} seconds")
        return ppt_buffer

    except Exception as e:
        log_message(logging.ERROR, f"Error in PowerPoint report generation: {str(e)}")
        raise
    finally:
        cleanup_images(user_id)

def generate_monotonicity_pdf_report(plots_data, user_id, well_name=None, model_equation=None):
    """
    Generate a PDF report containing monotonicity plots.
    
    :param plots_data: Dictionary containing stage numbers and their corresponding plotly figures
    :param user_id: str, user identifier
    :param well_name: str, name of the well being analyzed
    :param model_equation: str, the equation used in the model
    :return: BytesIO buffer containing the PDF
    """
    overall_start_time = time.time()
    log_message(logging.INFO, "Starting Monotonicity PDF report generation")
    buffer = BytesIO()
    
    try:
        doc = SimpleDocTemplate(buffer, pagesize=landscape(letter), 
                              rightMargin=36, leftMargin=36, 
                              topMargin=36, bottomMargin=18)
        
        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(
            name='CoverTitle',
            parent=styles['Title'],
            fontSize=24,
            spaceAfter=30
        ))
        styles.add(ParagraphStyle(
            name='WellName',
            parent=styles['Normal'],
            fontSize=16,
            spaceAfter=20
        ))
        styles.add(ParagraphStyle(
            name='ModelEquation',
            parent=styles['Normal'],
            fontSize=12,
            leftIndent=36,
            rightIndent=36,
            spaceAfter=30
        ))
        styles.add(ParagraphStyle(
            name='KeyAttributeAlert',
            parent=styles['Normal'],
            fontSize=14,
            textColor=colors.red,
            spaceAfter=10
        ))
        styles.add(ParagraphStyle(
            name='KeyAttributeValid',
            parent=styles['Normal'],
            fontSize=14,
            textColor=colors.green,
            spaceAfter=10
        ))
        
        elements = []
        
        # Cover Page
        elements.append(Paragraph("Monotonicity Analysis Report", styles['CoverTitle']))
        elements.append(Spacer(1, 20))
        
        # Add description of what monotonicity means
        elements.append(Paragraph("Monotonicity Analysis Overview:", styles['Heading2']))
        elements.append(Paragraph(
            "This report analyzes the monotonic behavior of the model - whether productivity consistently increases as key attributes increase. "
            "For hydraulic fracturing, we expect that increasing certain key parameters should lead to increased productivity.", 
            styles['Normal']))
        elements.append(Paragraph(
            "Key attributes of interest are:", 
            styles['Normal']))
        elements.append(Paragraph("• downhole_ppm - Downhole Proppant Concentration", styles['Normal']))
        elements.append(Paragraph("• total_dhppm - Total Downhole Proppant per Minute", styles['Normal']))
        elements.append(Paragraph("• tee - Treating Effective Energy", styles['Normal']))
        elements.append(Spacer(1, 12))
        
        if well_name:
            elements.append(Paragraph(f"Well Name: {well_name}", styles['WellName']))
        
        if model_equation:
            elements.append(Paragraph("Model Equation:", styles['Heading2']))
            # Split equation into multiple lines for better readability
            max_chars_per_line = 100
            equation_lines = []
            current_line = ""
            for word in model_equation.split():
                if len(current_line) + len(word) + 1 <= max_chars_per_line:
                    current_line += " " + word if current_line else word
                else:
                    equation_lines.append(current_line)
                    current_line = word
            if current_line:
                equation_lines.append(current_line)
            
            for line in equation_lines:
                elements.append(Paragraph(line, styles['ModelEquation']))
        
        # Add date and time of report generation
        current_time = time.strftime("%Y-%m-%d %H:%M:%S")
        elements.append(Paragraph(f"Generated on: {current_time}", styles['Normal']))
        
        elements.append(PageBreak())
        
        # Table of Contents
        elements.append(Paragraph("Contents", styles['Heading1']))
        elements.append(Spacer(1, 12))
        toc_items = [f"Stage {stage}" for stage in plots_data.keys()]
        for item in toc_items:
            elements.append(Paragraph(f"• {item}", styles['Normal']))
        
        elements.append(PageBreak())
        
        # Add plots for each stage
        for stage, plot_data in plots_data.items():
            elements.append(Paragraph(f"Stage {stage}", styles['Heading2']))
            elements.append(Spacer(1, 12))
            
            # Check for key attribute monotonicity data if available in the result_df object
            # The plot_data is the Figure, but we need to check if there's raw data available
            if hasattr(plot_data, '_data_obj') and hasattr(plot_data._data_obj, 'result_df'):
                result_df = plot_data._data_obj.result_df
                key_attributes = ['downhole_ppm', 'total_dhppm', 'tee']
                
                # Create a table for key attributes monotonicity
                monotonicity_data = []
                header_row = ['Attribute', 'Monotonicity %', 'Direction', 'Status']
                monotonicity_data.append(header_row)
                
                for attr in key_attributes:
                    pct_col = f"{attr}_monotonicity_pct"
                    dir_col = f"{attr}_monotonicity_dir"
                    
                    if pct_col in result_df.columns and dir_col in result_df.columns:
                        # Get the first non-null value
                        pct_values = result_df[pct_col].dropna()
                        dir_values = result_df[dir_col].dropna()
                        
                        if not pct_values.empty and not dir_values.empty:
                            pct = pct_values.iloc[0]
                            direction = dir_values.iloc[0]
                            status = "✓ VALID" if direction == "increasing" and pct >= 75 else "✗ INVALID"
                            
                            monotonicity_data.append([
                                attr,
                                f"{pct:.1f}%",
                                direction,
                                status
                            ])
                
                if len(monotonicity_data) > 1:  # If we have data beyond the header
                    # Create a table with the data
                    table = Table(monotonicity_data)
                    # Style the table
                    table_style = TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black)
                    ])
                    
                    # Add additional styling for valid/invalid status
                    for i, row in enumerate(monotonicity_data[1:], 1):
                        if "VALID" in row[-1]:
                            table_style.add('TEXTCOLOR', (3, i), (3, i), colors.green)
                        else:
                            table_style.add('TEXTCOLOR', (3, i), (3, i), colors.red)
                    
                    table.setStyle(table_style)
                    elements.append(table)
                    elements.append(Spacer(1, 12))
                
                # Add summary of monotonicity check
                valid_count = sum(1 for row in monotonicity_data[1:] if "VALID" in row[-1])
                total_count = len(monotonicity_data) - 1  # Subtract header row
                
                if valid_count == total_count:
                    elements.append(Paragraph("✅ All key attributes show proper increasing monotonic behavior.", styles['KeyAttributeValid']))
                else:
                    elements.append(Paragraph(f"⚠️ Only {valid_count} out of {total_count} key attributes show proper increasing monotonic behavior.", styles['KeyAttributeAlert']))
                
                elements.append(Spacer(1, 12))
            
            try:
                img_path = save_plot_as_image(plot_data, f"monotonicity_stage_{stage}", user_id)
                img = Image(img_path, width=10*inch, height=5.6*inch)
                elements.append(img)
            except Exception as e:
                log_message(logging.ERROR, f"Error adding plot for Stage {stage} to PDF: {str(e)}", exc_info=True)
                elements.append(Paragraph(f"Error including plot for Stage {stage}: {str(e)}", styles['Normal']))
            
            # Add individual plots for key attributes if available
            if hasattr(plot_data, '_data_obj') and hasattr(plot_data._data_obj, 'key_attr_plots'):
                key_attr_plots = plot_data._data_obj.key_attr_plots
                
                for attr, attr_plot in key_attr_plots.items():
                    try:
                        elements.append(Paragraph(f"Monotonicity Analysis for {attr}", styles['Heading3']))
                        img_path = save_plot_as_image(attr_plot, f"monotonicity_{attr}_stage_{stage}", user_id)
                        img = Image(img_path, width=10*inch, height=5.6*inch)
                        elements.append(img)
                    except Exception as e:
                        log_message(logging.ERROR, f"Error adding {attr} plot for Stage {stage} to PDF: {str(e)}", exc_info=True)
            
            elements.append(Spacer(1, 12))
            elements.append(PageBreak())
        
        # Add a summary page at the end
        elements.append(Paragraph("Summary of Monotonicity Analysis", styles['Heading1']))
        elements.append(Spacer(1, 12))
        
        elements.append(Paragraph(
            "Monotonicity Check Interpretation:", 
            styles['Heading3']))
        
        elements.append(Paragraph(
            "• A monotonic relationship means productivity should consistently increase (or consistently decrease) as an attribute increases.", 
            styles['Normal']))
        
        elements.append(Paragraph(
            "• For hydraulic fracturing models, we specifically want to see that productivity increases as key attributes increase.", 
            styles['Normal']))
        
        elements.append(Paragraph(
            "• A model with proper monotonicity is more physically realistic and reliable for operational decisions.", 
            styles['Normal']))
        
        elements.append(Paragraph(
            "• Non-monotonic behavior may indicate issues with the model or unexpected physical phenomena that should be investigated.", 
            styles['Normal']))
        
        # Build the PDF
        log_message(logging.DEBUG, "Building PDF")
        pdf_start_time = time.time()
        doc.build(elements)
        log_message(logging.DEBUG, f"Built PDF in {time.time() - pdf_start_time:.2f} seconds")
        
        # Reset buffer position
        buffer.seek(0)
        
        overall_end_time = time.time()
        log_message(logging.INFO, f"Monotonicity PDF report generated successfully in {overall_end_time - overall_start_time:.2f} seconds")
        return buffer
        
    except Exception as e:
        log_message(logging.ERROR, f"Error in Monotonicity PDF report generation: {str(e)}", exc_info=True)
        raise
    finally:
        if 'doc' in locals():
            del doc
        cleanup_images(user_id)

if __name__ == "__main__":
    # You can add any test code here if needed
    pass